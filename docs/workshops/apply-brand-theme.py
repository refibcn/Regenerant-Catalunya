#!/usr/bin/env python3
"""
Apply Regenerant Catalunya brand theme to PowerPoint presentation.

This script takes a pandoc-generated PPTX file and applies:
- Brand colors (Forest Green, Warm Sand, etc.)
- Typography (Inter font family)
- Design principles (straight edges, bold borders)
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Regenerant Catalunya Brand Colors
COLORS = {
    'forest_green': RGBColor(33, 64, 51),      # #214033
    'warm_sand': RGBColor(230, 223, 215),     # #E6DFD7
    'midnight_navy': RGBColor(9, 32, 69),      # #092045
    'green_medium': RGBColor(58, 102, 85),     # #3A6655
    'green_light': RGBColor(86, 143, 121),     # #568F79
    'purple': RGBColor(107, 78, 163),          # #6B4EA3
    'orange': RGBColor(230, 126, 80),          # #E67E50
}

# Font settings
FONT_FAMILY = 'Inter'
FONT_FALLBACK = 'Arial'  # Fallback if Inter not available

def hex_to_rgb(hex_color):
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def apply_slide_background(slide, color):
    """Apply background color to slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def apply_text_formatting(text_frame, is_header=False):
    """Apply typography to text frame."""
    for paragraph in text_frame.paragraphs:
        # Set font
        for run in paragraph.runs:
            # Try Inter first, fallback to Arial
            try:
                run.font.name = FONT_FAMILY
            except:
                run.font.name = FONT_FALLBACK
            
            if is_header:
                run.font.bold = True
                run.font.size = Pt(32) if paragraph.level == 0 else Pt(24)
                run.font.color.rgb = COLORS['forest_green']
            else:
                run.font.size = Pt(18)
                run.font.color.rgb = COLORS['midnight_navy']
        
        # Set paragraph spacing
        paragraph.space_after = Pt(12)
        paragraph.space_before = Pt(6)

def style_title_slide(slide):
    """Style title slide with brand colors."""
    apply_slide_background(slide, COLORS['warm_sand'])
    
    # Style title
    if slide.shapes.title:
        title = slide.shapes.title
        try:
            title.text_frame.paragraphs[0].font.name = FONT_FAMILY
        except:
            title.text_frame.paragraphs[0].font.name = FONT_FALLBACK
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLORS['forest_green']
    
    # Style subtitle if exists
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
            apply_text_formatting(shape.text_frame, is_header=False)

def style_content_slide(slide):
    """Style content slide with brand colors."""
    apply_slide_background(slide, COLORS['warm_sand'])
    
    # Style title
    if slide.shapes.title:
        title = slide.shapes.title
        try:
            title.text_frame.paragraphs[0].font.name = FONT_FAMILY
        except:
            title.text_frame.paragraphs[0].font.name = FONT_FALLBACK
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLORS['forest_green']
        
        # Add border line under title
        # Note: python-pptx doesn't directly support borders on text boxes,
        # but we can add a line shape below the title
    
    # Style content
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
            apply_text_formatting(shape.text_frame, is_header=False)

def add_title_border(slide):
    """Add a bold border line under the title."""
    if slide.shapes.title:
        title = slide.shapes.title
        title_bottom = title.top + title.height
        
        # Add a line shape (2-3px thick)
        line = slide.shapes.add_connector(1, Inches(0.5), title_bottom + Pt(10), 
                                         Inches(9.5), title_bottom + Pt(10))
        line.line.color.rgb = COLORS['forest_green']
        line.line.width = Pt(3)

def process_presentation(input_file, output_file):
    """Process PowerPoint file and apply brand theme."""
    print(f"Loading presentation: {input_file}")
    prs = Presentation(input_file)
    
    print(f"Processing {len(prs.slides)} slides...")
    
    for i, slide in enumerate(prs.slides):
        # Determine slide type
        if i == 0 or (slide.shapes.title and 
                     ('title' in slide.shapes.title.text.lower() or 
                      'workshop' in slide.shapes.title.text.lower())):
            style_title_slide(slide)
        else:
            style_content_slide(slide)
            # Try to add border under title
            try:
                add_title_border(slide)
            except:
                pass  # Skip if we can't add border
    
    print(f"Saving styled presentation: {output_file}")
    prs.save(output_file)
    print("✓ Presentation styled successfully!")

def create_theme_template(output_file):
    """Create a reusable theme template with master slides."""
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Title slide layout
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_slide_background(title_slide, COLORS['warm_sand'])
    if title_slide.shapes.title:
        title_slide.shapes.title.text = "Title Slide"
        style_title_slide(title_slide)
    
    # Content slide layout
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_background(content_slide, COLORS['warm_sand'])
    if content_slide.shapes.title:
        content_slide.shapes.title.text = "Content Slide"
        style_content_slide(content_slide)
    
    # Two-column layout
    two_col_slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_background(two_col_slide, COLORS['warm_sand'])
    if two_col_slide.shapes.title:
        two_col_slide.shapes.title.text = "Two Column Layout"
        style_content_slide(two_col_slide)
    
    print(f"Saving theme template: {output_file}")
    prs.save(output_file)
    print("✓ Theme template created successfully!")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage:")
        print("  Apply theme to presentation:")
        print("    python3 apply-brand-theme.py <input.pptx> <output.pptx>")
        print("  Create theme template:")
        print("    python3 apply-brand-theme.py --theme <output.pptx>")
        sys.exit(1)
    
    if sys.argv[1] == '--theme':
        output_file = sys.argv[2] if len(sys.argv) > 2 else 'regenerant-catalunya-theme.pptx'
        create_theme_template(output_file)
    else:
        input_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else input_file.replace('.pptx', '-styled.pptx')
        process_presentation(input_file, output_file)

